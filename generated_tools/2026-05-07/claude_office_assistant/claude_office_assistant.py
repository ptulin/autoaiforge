import argparse
import os
import openai
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

def summarize_word(file_path, api_key):
    """Summarize the content of a Word document using Claude AI."""
    if not os.path.exists(file_path):
        raise FileNotFoundError("The specified Word file does not exist.")

    document = Document(file_path)
    text = "\n".join([paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()])

    if not text.strip():
        return "The document is empty."

    openai.api_key = api_key

    try:
        response = openai.Completion.create(
            engine="text-davinci-003",
            prompt=f"Summarize the following document:\n{text}",
            max_tokens=200
        )
        return response.choices[0].text.strip()
    except Exception as e:
        return f"Error during summarization: {str(e)}"

def analyze_excel(file_path):
    """Analyze an Excel file and return basic statistics."""
    if not os.path.exists(file_path):
        raise FileNotFoundError("The specified Excel file does not exist.")

    workbook = load_workbook(file_path)
    analysis = {}

    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        rows = list(worksheet.iter_rows(values_only=True))
        analysis[sheet] = {
            "row_count": len(rows),
            "column_count": len(rows[0]) if rows else 0
        }

    return analysis

def generate_presentation(text, output_file):
    """Generate a PowerPoint presentation from the provided text."""
    if not text.strip():
        raise ValueError("Input text is empty.")

    presentation = Presentation()
    slides = text.split("\n\n")

    for slide_text in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title, *content = slide_text.split("\n")
        slide.shapes.title.text = title
        slide.placeholders[1].text = "\n".join(content)

    presentation.save(output_file)
    return f"Presentation saved to {output_file}"

def main():
    parser = argparse.ArgumentParser(description="Claude Office Assistant")
    parser.add_argument("--file", required=True, help="Path to the input file (Word, Excel, or PowerPoint).")
    parser.add_argument("--summarize", action="store_true", help="Summarize the content of a Word document.")
    parser.add_argument("--analyze", action="store_true", help="Analyze the content of an Excel file.")
    parser.add_argument("--generate", help="Generate a PowerPoint presentation from text.")
    parser.add_argument("--output", help="Path to save the output file.")
    parser.add_argument("--api_key", help="API key for Claude AI.")

    args = parser.parse_args()

    if args.summarize:
        if not args.api_key:
            print("Error: --api_key is required for summarization.")
            return
        summary = summarize_word(args.file, args.api_key)
        print(summary)

    elif args.analyze:
        analysis = analyze_excel(args.file)
        print(analysis)

    elif args.generate:
        if not args.output:
            print("Error: --output is required for generating a presentation.")
            return
        result = generate_presentation(args.generate, args.output)
        print(result)

    else:
        print("Error: No valid operation specified. Use --summarize, --analyze, or --generate.")

if __name__ == "__main__":
    main()